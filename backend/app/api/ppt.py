"""PPT 联动与知识点检索：列表、按章节/关键词检索、教师上传"""
import os
import re
import time
from pathlib import Path

from fastapi import APIRouter, Depends, Query, UploadFile, File, Form, HTTPException
from pydantic import BaseModel
from sqlalchemy import select, or_
from sqlalchemy.ext.asyncio import AsyncSession
from pptx import Presentation

from ..config import settings
from ..db import get_db
from ..db.models import PptFile, PptSlide, Chapter
from ..api.auth import require_teacher
from ..db.models import User

router = APIRouter(prefix="/ppt", tags=["ppt"])


class PptFileOut(BaseModel):
    id: int
    chapter_id: int | None
    file_name: str
    total_slides: int

    class Config:
        from_attributes = True


class PptSlideOut(BaseModel):
    id: int
    ppt_id: int
    slide_index: int
    text_content: str | None
    keywords: str | None

    class Config:
        from_attributes = True


class PptSlideSearchOut(PptSlideOut):
    chapter_title: str | None = None
    file_name: str | None = None


@router.get("", response_model=list[PptFileOut])
async def list_ppts(
    chapter_id: int | None = Query(None),
    db: AsyncSession = Depends(get_db),
):
    """按章节列出 PPT 文件（可选 chapter_id 筛选）"""
    q = select(PptFile).order_by(PptFile.chapter_id, PptFile.id)
    if chapter_id is not None:
        q = q.where(PptFile.chapter_id == chapter_id)
    r = await db.execute(q)
    return [PptFileOut.model_validate(x) for x in r.scalars().all()]


@router.get("/search/slides", response_model=list[PptSlideSearchOut])
async def search_slides(
    chapter_id: int | None = Query(None, description="按章节筛选"),
    keyword: str | None = Query(None, description="关键词，匹配 text_content 或 keywords"),
    limit: int = Query(20, le=50),
    db: AsyncSession = Depends(get_db),
):
    """按章节、关键词检索幻灯片，用于课中 PPT 定位与知识点跳转"""
    q = (
        select(PptSlide, PptFile, Chapter.title.label("chapter_title"))
        .join(PptFile, PptSlide.ppt_id == PptFile.id)
        .outerjoin(Chapter, PptFile.chapter_id == Chapter.id)
    )
    if chapter_id is not None:
        q = q.where(PptFile.chapter_id == chapter_id)
    if keyword and keyword.strip():
        k = f"%{keyword.strip()[:64]}%"
        q = q.where(
            or_(
                PptSlide.text_content.ilike(k),
                PptSlide.keywords.ilike(k),
            )
        )
    q = q.order_by(PptFile.chapter_id, PptSlide.slide_index).limit(limit)
    r = await db.execute(q)
    rows = r.all()
    return [
        PptSlideSearchOut(
            id=row[0].id,
            ppt_id=row[0].ppt_id,
            slide_index=row[0].slide_index,
            text_content=row[0].text_content,
            keywords=row[0].keywords,
            chapter_title=row[2],
            file_name=row[1].file_name if row[1] else None,
        )
        for row in rows
    ]


@router.get("/{ppt_id}/slides", response_model=list[PptSlideOut])
async def list_slides(
    ppt_id: int,
    db: AsyncSession = Depends(get_db),
):
    """获取指定 PPT 的幻灯片列表（页码、文本、关键词）"""
    r = await db.execute(select(PptFile).where(PptFile.id == ppt_id))
    if not r.scalar_one_or_none():
        raise HTTPException(status_code=404, detail="PPT 不存在")
    r = await db.execute(
        select(PptSlide).where(PptSlide.ppt_id == ppt_id).order_by(PptSlide.slide_index)
    )
    return [PptSlideOut.model_validate(x) for x in r.scalars().all()]


def _safe_filename(name: str) -> str:
    """保留扩展名，仅允许字母数字、横线、下划线、点"""
    base, ext = os.path.splitext(name)
    safe = re.sub(r"[^\w\-.]", "_", base)[:80]
    return (safe or "ppt") + (ext.lower() if ext.lower() == ".pptx" else ".pptx")


@router.post("/upload", response_model=PptFileOut)
async def upload_ppt(
    chapter_id: int = Form(...),
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    user: User = Depends(require_teacher),
):
    """教师/管理员上传 PPTX，按章节关联；解析每页备注与正文用于检索"""
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="请上传 .pptx 文件")
    r = await db.execute(select(Chapter).where(Chapter.id == chapter_id))
    chapter = r.scalar_one_or_none()
    if not chapter:
        raise HTTPException(status_code=404, detail="章节不存在")

    root = Path(settings.upload_dir)
    root.mkdir(parents=True, exist_ok=True)
    subdir = root / "ppt"
    subdir.mkdir(parents=True, exist_ok=True)
    safe_name = _safe_filename(file.filename)
    stem = f"{chapter_id}_{int(time.time())}_{safe_name}"
    file_path = subdir / stem
    content = await file.read()
    file_path.write_bytes(content)

    try:
        prs = Presentation(file_path)
    except Exception as e:
        file_path.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"PPT 解析失败: {str(e)}")

    total_slides = len(prs.slides)
    rel_path = f"ppt/{stem}"
    ppt_entity = PptFile(
        chapter_id=chapter_id,
        file_name=file.filename or safe_name,
        file_path=rel_path,
        total_slides=total_slides,
    )
    db.add(ppt_entity)
    await db.flush()

    for idx, slide in enumerate(prs.slides):
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            pass
        shape_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                shape_texts.append(shape.text.strip())
        body_text = "\n".join(shape_texts).strip()
        text_content = (notes_text + "\n" + body_text).strip() or None
        db.add(PptSlide(ppt_id=ppt_entity.id, slide_index=idx + 1, text_content=text_content))
    await db.commit()
    await db.refresh(ppt_entity)
    return PptFileOut.model_validate(ppt_entity)
