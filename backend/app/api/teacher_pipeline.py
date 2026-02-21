"""教师端课程素材流水线：PDF -> PPT -> 讲解脚本 -> 音频 -> 视频。

该模块按阶段落盘，支持逐段执行和人工编辑后继续执行。
"""
from __future__ import annotations

import hashlib
import json
import io
import os
import re
import base64
import logging
import html
import time
import shutil
import subprocess
import tempfile
import textwrap
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any
from urllib.parse import quote

from fastapi import APIRouter, Depends, File, Form, HTTPException, Query, UploadFile
from fastapi.responses import FileResponse, Response
from openai import OpenAI
from pydantic import BaseModel, Field
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Pt
import requests

from ..api.auth import require_teacher
from ..config import settings
from ..db.models import User

router = APIRouter(prefix="/teacher/pipeline", tags=["teacher-pipeline"])
logger = logging.getLogger(__name__)

VOICE_ZH_LABELS = {
    "Serena": "赛瑞娜",
    "Ethan": "伊森",
    "Chelsie": "切尔西",
    "Momo": "莫莫",
    "Vivian": "薇薇安",
    "Moon": "慕恩",
    "Maia": "玛雅",
    "Kai": "凯",
    "Nofish": "小鱼",
    "Bella": "贝拉",
    "Jennifer": "珍妮弗",
    "Ryan": "瑞安",
    "Katerina": "卡特里娜",
    "Aiden": "艾登",
    "Cherry": "樱桃",
    "Ono Anna": "小野杏",
    "Lenn": "莱恩",
    "Emilien": "埃米尔安",
    "Andre": "安德雷",
    "Radio Gol": "拉迪奥·戈尔",
    "Jada": "上海-阿珍",
    "Dylan": "北京-晓东",
    "Li": "南京-老李",
    "Marcus": "陕西-秦川",
    "Roy": "闽南-阿杰",
    "Peter": "天津-李彼得",
    "Sunny": "四川-晴儿",
    "Eric": "四川-程川",
    "Rocky": "粤语-阿强",
    "Kiki": "粤语-阿清",
}

# 音色性别：用于前端按性别筛选。未知音色不包含在 map 中，前端显示为「全部」可选。
VOICE_GENDER: dict[str, str] = {
    # 千问 qwen3-tts-flash 常见音色（女声）
    "Serena": "female",
    "Chelsie": "female",
    "Momo": "female",
    "Vivian": "female",
    "Moon": "female",
    "Maia": "female",
    "Bella": "female",
    "Jennifer": "female",
    "Katerina": "female",
    "Cherry": "female",
    "Ono Anna": "female",
    "Jada": "female",
    "Sunny": "female",
    "Kiki": "female",
    # 千问 男声
    "Ethan": "male",
    "Kai": "male",
    "Ryan": "male",
    "Aiden": "male",
    "Lenn": "male",
    "Emilien": "male",
    "Andre": "male",
    "Dylan": "male",
    "Li": "male",
    "Marcus": "male",
    "Roy": "male",
    "Peter": "male",
    "Eric": "male",
    "Rocky": "male",
    "Nofish": "male",
    "Radio Gol": "male",
    # OpenAI 兼容 TTS 音色
    "alloy": "female",
    "echo": "male",
    "fable": "female",
    "onyx": "male",
    "nova": "female",
    "shimmer": "female",
}

_QWEN_DOC_VOICE_CACHE: dict[str, Any] = {"at": 0.0, "map": {}}


def _now_iso() -> str:
    return datetime.utcnow().replace(microsecond=0).isoformat() + "Z"


def _pipeline_root() -> Path:
    root = Path(settings.upload_dir) / "teacher_pipeline"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _safe_name(name: str) -> str:
    return re.sub(r"[^\w\-.]", "_", name).strip("._")[:120] or "unnamed"


def _workflow_dir(workflow_id: str) -> Path:
    if not re.fullmatch(r"[a-zA-Z0-9_-]{3,64}", workflow_id):
        raise HTTPException(status_code=400, detail="workflow_id 不合法")
    path = _pipeline_root() / workflow_id
    if not path.exists():
        raise HTTPException(status_code=404, detail="工作流不存在")
    return path


def _manifest_path(workflow_id: str) -> Path:
    return _workflow_dir(workflow_id) / "manifest.json"


def _read_manifest(workflow_id: str) -> dict[str, Any]:
    path = _manifest_path(workflow_id)
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"manifest 读取失败: {str(e)}")


def _write_manifest(workflow_id: str, manifest: dict[str, Any]) -> None:
    path = _manifest_path(workflow_id)
    path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def _touch_stage(
    workflow_id: str,
    stage: str,
    status: str,
    outputs: list[str] | None = None,
    extra: dict[str, Any] | None = None,
) -> dict[str, Any]:
    manifest = _read_manifest(workflow_id)
    stages = manifest.setdefault("stages", {})
    stage_data = stages.setdefault(stage, {})
    stage_data["status"] = status
    stage_data["updated_at"] = _now_iso()
    if outputs is not None:
        stage_data["outputs"] = outputs
    if extra:
        stage_data.update(extra)
    manifest["updated_at"] = _now_iso()
    _write_manifest(workflow_id, manifest)
    return manifest


def _safe_workflow_file(workflow_id: str, rel_path: str) -> Path:
    wf_dir = _workflow_dir(workflow_id).resolve()
    rel = rel_path.strip().lstrip("/").replace("\\", "/")
    if not rel:
        raise HTTPException(status_code=400, detail="path 不能为空")
    full = (wf_dir / rel).resolve()
    if not str(full).startswith(str(wf_dir)):
        raise HTTPException(status_code=400, detail="非法路径")
    return full


def _extract_pdf_text(pdf_bytes: bytes) -> tuple[str, int]:
    try:
        reader = PdfReader(io.BytesIO(pdf_bytes))
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"PDF 解析失败: {str(e)}")

    pages: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        txt = (page.extract_text() or "").strip()
        if txt:
            pages.append(f"\n\n## 第{i}页\n{txt}")
    return ("\n".join(pages)).strip(), len(reader.pages)


def _split_text_by_chapter(text: str) -> list[dict[str, Any]]:
    lines = text.splitlines()
    chapter_pat = re.compile(r"^\s*(第[0-9一二三四五六七八九十百千]+章[^\n]{0,50})\s*$")
    page_pat = re.compile(r"^\s*##\s*第\d+页\s*$")
    indexes: list[tuple[int, str]] = []
    for i, line in enumerate(lines):
        if page_pat.match(line):
            continue
        m = chapter_pat.match(line)
        if m:
            indexes.append((i, m.group(1).strip()))

    if not indexes:
        whole = "\n".join(lines).strip()
        return [{"chapter_no": 1, "title": "全文", "content": whole}]

    chunks: list[dict[str, Any]] = []
    for idx, (start, title) in enumerate(indexes):
        end = indexes[idx + 1][0] if idx + 1 < len(indexes) else len(lines)
        content = "\n".join(lines[start:end]).strip()
        if not content:
            continue
        chunks.append({"chapter_no": len(chunks) + 1, "title": title, "content": content})
    return chunks or [{"chapter_no": 1, "title": "全文", "content": text.strip()}]


def _workflow_id_for_doc(doc_id: int) -> str:
    return f"doc_{doc_id}"


def _ensure_workflow_dirs(workflow_id: str) -> Path:
    wf_dir = _pipeline_root() / workflow_id
    wf_dir.mkdir(parents=True, exist_ok=True)
    for sub in ["stage1", "stage2", "stage3", "stage4", "stage5", "stage6"]:
        (wf_dir / sub).mkdir(parents=True, exist_ok=True)
    return wf_dir


def _docs_root() -> Path:
    root = _pipeline_root() / "docs"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _pdf_docs_index_path() -> Path:
    return _pipeline_root() / "pdf_docs_index.json"


def _read_pdf_docs_index() -> list[dict[str, Any]]:
    p = _pdf_docs_index_path()
    if not p.exists():
        return []
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _write_pdf_docs_index(rows: list[dict[str, Any]]) -> None:
    p = _pdf_docs_index_path()
    p.write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")


def _find_pdf_doc_or_404(user_id: int, doc_id: int) -> dict[str, Any]:
    for row in _read_pdf_docs_index():
        if int(row.get("id") or 0) == doc_id and int(row.get("owner_id") or 0) == user_id:
            return row
    raise HTTPException(status_code=404, detail="PDF 文档不存在或无权限")


def _chunk_text(text: str, chunk_size: int = 1200, chunk_overlap: int = 120) -> list[dict[str, Any]]:
    _ = (text, chunk_size, chunk_overlap)
    return []


def _fallback_build_slides(text: str, max_slides: int = 20) -> list[dict[str, Any]]:
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    if not paras:
        paras = [text.strip()] if text.strip() else []

    slides: list[dict[str, Any]] = []
    for i, para in enumerate(paras[:max_slides], start=1):
        lines = [x.strip(" -•\t") for x in para.split("\n") if x.strip()]
        title = (lines[0] if lines else para[:24]).strip()[:28] or f"教学要点{i}"
        bullets = []
        for line in lines[1:5]:
            cleaned = line[:60]
            if cleaned:
                bullets.append(cleaned)
        if not bullets:
            summary = para[:140]
            bullets = textwrap.wrap(summary, width=28)[:4]
        slides.append(
            {
                "slide_no": i,
                "title": title,
                "bullets": bullets,
                "notes": f"本页重点：{title}。请结合教材原文展开讲解。",
            }
        )
    return slides


def _openai_client() -> OpenAI:
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("未配置 OPENAI_API_KEY")
    return OpenAI(api_key=api_key)


def _resolve_tts_client_and_model(model: str) -> tuple[OpenAI, str]:
    model_in = (model or "").strip()
    if ":" in model_in:
        provider_id, real_model = model_in.split(":", 1)
        provider_id = provider_id.strip()
        real_model = real_model.strip()
        if not provider_id or not real_model:
            raise HTTPException(status_code=400, detail="TTS 模型格式不正确，应为 provider_id:model")
        try:
            from ..rag.config_store import get_providers_list_raw
            providers = get_providers_list_raw()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"无法读取模型提供商配置: {str(e)}")
        p = next((x for x in providers if (x.get("id") or "").strip() == provider_id), None)
        if not p:
            raise HTTPException(status_code=400, detail="未找到对应模型提供商，请联系管理员检查配置")
        api_key = (p.get("api_key") or "").strip()
        if not api_key:
            raise HTTPException(status_code=400, detail="该模型提供商未配置 API Key")
        provider_type = (p.get("type") or "").strip() or "openai_compatible"
        base_url = (p.get("base_url") or "").strip()
        if not base_url:
            # Compatible OpenAI endpoint defaults by provider.
            if provider_type == "qianwen":
                # Beijing endpoint. For international keys, set provider base_url to dashscope-intl endpoint in Admin.
                base_url = "https://dashscope.aliyuncs.com/compatible-mode/v1"
            elif provider_type == "zhipu":
                base_url = "https://open.bigmodel.cn/api/paas/v4"
        return OpenAI(api_key=api_key, base_url=base_url), real_model

    if not model_in:
        try:
            from ..rag.config_store import get_default_llm, get_default_tts
            default_tts = (get_default_tts() or "").strip()
            default_llm = (get_default_llm() or "").strip()
        except Exception:
            default_tts = ""
            default_llm = ""
        if ":" in default_tts:
            return _resolve_tts_client_and_model(default_tts)
        if ":" in default_llm:
            return _resolve_tts_client_and_model(default_llm)
        model_in = "gpt-4o-mini-tts"

    return _openai_client(), model_in


def _resolve_provider_and_model(model: str) -> tuple[dict[str, Any] | None, str]:
    model_in = (model or "").strip()
    if ":" not in model_in:
        return None, model_in
    provider_id, real_model = model_in.split(":", 1)
    provider_id = provider_id.strip()
    real_model = real_model.strip()
    if not provider_id or not real_model:
        raise HTTPException(status_code=400, detail="TTS 模型格式不正确，应为 provider_id:model")
    try:
        from ..rag.config_store import get_providers_list_raw
        providers = get_providers_list_raw()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"无法读取模型提供商配置: {str(e)}")
    p = next((x for x in providers if (x.get("id") or "").strip() == provider_id), None)
    if not p:
        raise HTTPException(status_code=400, detail="未找到对应模型提供商，请联系管理员检查配置")
    return p, real_model


def _normalize_qianwen_api_root(base_url: str) -> str:
    b = (base_url or "").strip().rstrip("/")
    if not b:
        return "https://dashscope.aliyuncs.com/api/v1"
    if "dashscope-intl.aliyuncs.com" in b:
        return "https://dashscope-intl.aliyuncs.com/api/v1"
    if "dashscope.aliyuncs.com" in b:
        return "https://dashscope.aliyuncs.com/api/v1"
    if b.endswith("/api/v1"):
        return b
    if b.endswith("/v1"):
        return b[:-3] + "/api/v1"
    return b + "/api/v1"


def _collect_string_values_by_keys(data: Any, key_predicate) -> list[str]:
    out: list[str] = []

    def walk(node: Any):
        if isinstance(node, dict):
            for k, v in node.items():
                if key_predicate(str(k)):
                    if isinstance(v, str):
                        out.append(v.strip())
                    elif isinstance(v, list):
                        for item in v:
                            if isinstance(item, str):
                                out.append(item.strip())
                            elif isinstance(item, dict):
                                for kk in ("id", "name", "value", "voice", "model"):
                                    vv = item.get(kk)
                                    if isinstance(vv, str):
                                        out.append(vv.strip())
                    elif isinstance(v, dict):
                        for kk in ("id", "name", "value", "voice", "model"):
                            vv = v.get(kk)
                            if isinstance(vv, str):
                                out.append(vv.strip())
                walk(v)
        elif isinstance(node, list):
            for item in node:
                walk(item)

    walk(data)
    return [x for x in out if x]


def _extract_tts_model_ids(data: Any) -> list[str]:
    # Try direct keys first.
    direct = _collect_string_values_by_keys(
        data,
        lambda k: k.lower() in {"id", "model", "model_id", "name"},
    )
    cand = []
    for m in direct:
        mm = m.strip()
        low = mm.lower()
        if "tts" in low and len(mm) <= 128:
            cand.append(mm)
    # Preserve order and remove duplicates.
    seen = set()
    out: list[str] = []
    for x in cand:
        if x not in seen:
            seen.add(x)
            out.append(x)
    return out


def _extract_voices_from_model_meta(data: Any) -> list[str]:
    vals = _collect_string_values_by_keys(
        data,
        lambda k: ("voice" in k.lower()) or (k.lower() in {"speaker", "speakers"}),
    )
    voices: list[str] = []
    for v in vals:
        vv = v.strip()
        if not vv or len(vv) > 80:
            continue
        # Skip non-voice marker strings.
        if vv.lower() in {"voice", "voices", "model", "audio", "wav", "mp3"}:
            continue
        voices.append(vv)
    seen = set()
    out: list[str] = []
    for x in voices:
        if x not in seen:
            seen.add(x)
            out.append(x)
    return out


def _voice_label_zh(voice: str) -> str:
    cn = VOICE_ZH_LABELS.get(voice)
    return f"{cn}（{voice}）" if cn else voice


def _voice_gender(voice: str) -> str | None:
    """返回 'male' | 'female'，未知则返回 None。"""
    if not voice:
        return None
    return VOICE_GENDER.get(voice.strip())


def _normalize_model_family(model_id: str) -> str:
    m = (model_id or "").strip().lower()
    if not m:
        return m
    # remove date snapshot suffix
    m = re.sub(r"-\d{4}-\d{2}-\d{2}$", "", m)
    return m


def _fetch_qwen_doc_voice_matrix() -> dict[str, set[str]]:
    now = time.time()
    if now - float(_QWEN_DOC_VOICE_CACHE.get("at") or 0) < 3600 and _QWEN_DOC_VOICE_CACHE.get("map"):
        return _QWEN_DOC_VOICE_CACHE["map"]
    url = "https://help.aliyun.com/zh/model-studio/qwen-tts"
    logger.warning("[TTS][QWEN] fetch doc voice matrix url=%s", url)
    resp = requests.get(url, timeout=30)
    if resp.status_code >= 400:
        raise RuntimeError(f"拉取百炼文档失败 HTTP {resp.status_code}")
    text = resp.text
    matrix: dict[str, set[str]] = {}
    rows = re.findall(r"<tr[^>]*>(.*?)</tr>", text, flags=re.S | re.I)
    for row in rows:
        vm = re.search(r'class="code">([^<]{1,80})</code>', row, flags=re.I)
        if not vm:
            continue
        voice = html.unescape(vm.group(1)).strip()
        if not voice:
            continue
        model_ids = re.findall(
            r"\bqwen(?:3)?-tts(?:-[a-z0-9]+)*(?:-\d{4}-\d{2}-\d{2})?\b",
            row,
            flags=re.I,
        )
        for mid in model_ids:
            key = _normalize_model_family(mid)
            if not key:
                continue
            matrix.setdefault(key, set()).add(voice)
    # fallback minimal set if parse misses
    if "qwen3-tts-flash" not in matrix:
        matrix["qwen3-tts-flash"] = set(
            [
                "Serena",
                "Ethan",
                "Chelsie",
                "Momo",
                "Vivian",
                "Moon",
                "Maia",
                "Kai",
                "Nofish",
                "Bella",
                "Jennifer",
                "Ryan",
                "Katerina",
                "Aiden",
                "Cherry",
                "Ono Anna",
                "Lenn",
                "Emilien",
                "Andre",
                "Radio Gol",
                "Jada",
                "Dylan",
                "Li",
                "Marcus",
                "Roy",
                "Peter",
                "Sunny",
                "Eric",
                "Rocky",
                "Kiki",
            ]
        )
    _QWEN_DOC_VOICE_CACHE["at"] = now
    _QWEN_DOC_VOICE_CACHE["map"] = matrix
    logger.warning("[TTS][QWEN] doc voice matrix models=%s", len(matrix))
    return matrix


def _voices_allowed_by_doc(model_id: str) -> set[str]:
    family = _normalize_model_family(model_id)
    matrix = _fetch_qwen_doc_voice_matrix()
    return matrix.get(family, set())


def _fetch_qwen_doc_tts_models() -> list[str]:
    url = "https://help.aliyun.com/zh/model-studio/qwen-tts"
    logger.warning("[TTS][QWEN] fetch doc tts models url=%s", url)
    resp = requests.get(url, timeout=30)
    if resp.status_code >= 400:
        raise RuntimeError(f"拉取百炼文档模型失败 HTTP {resp.status_code}")
    text = resp.text
    # Extract model IDs from docs; exclude realtime variants for offline stage5 audio synthesis.
    ms = re.findall(
        r"\bqwen(?:3)?-tts(?:-[a-z0-9]+)*(?:-\d{4}-\d{2}-\d{2})?\b",
        text,
        flags=re.I,
    )
    out: list[str] = []
    seen = set()
    for m in ms:
        mm = m.strip()
        if not mm or "realtime" in mm.lower():
            continue
        if "voice-cloning" in mm.lower() or "voice-design" in mm.lower() or mm.lower().endswith("-api"):
            continue
        if mm in seen:
            continue
        seen.add(mm)
        out.append(mm)
    logger.warning("[TTS][QWEN] doc tts models count=%s", len(out))
    return out


def _fetch_qianwen_tts_catalog(api_key: str, base_url: str) -> tuple[list[str], dict[str, list[str]]]:
    root = _normalize_qianwen_api_root(base_url)
    headers = {"Authorization": f"Bearer {api_key}"}
    models_url = root.rstrip("/") + "/models"
    logger.warning("[TTS][QWEN] catalog fetch models_url=%s", models_url)
    resp = requests.get(models_url, headers=headers, timeout=30)
    logger.warning("[TTS][QWEN] catalog models http_status=%s", resp.status_code)
    if resp.status_code >= 400:
        raise RuntimeError(f"获取模型列表失败 HTTP {resp.status_code}: {resp.text[:300]}")
    payload = resp.json()
    logger.warning(
        "[TTS][QWEN] catalog models payload_type=%s keys=%s sample=%s",
        type(payload).__name__,
        list(payload.keys())[:20] if isinstance(payload, dict) else [],
        json.dumps(payload, ensure_ascii=False)[:400],
    )
    model_ids = _extract_tts_model_ids(payload)
    if not model_ids:
        # Fallback: regex extraction from raw payload text.
        raw = json.dumps(payload, ensure_ascii=False)
        ms = re.findall(
            r"\bqwen(?:3)?-tts(?:-[a-z0-9]+)*(?:-\d{4}-\d{2}-\d{2})?\b",
            raw,
            flags=re.I,
        )
        seen = set()
        for m in ms:
            mm = m.strip()
            if not mm or "realtime" in mm.lower():
                continue
            if mm in seen:
                continue
            seen.add(mm)
            model_ids.append(mm)
        logger.warning("[TTS][QWEN] regex fallback models_count=%s", len(model_ids))
    if not model_ids:
        # Final fallback: dynamic extraction from official docs (not hardcoded).
        try:
            model_ids = _fetch_qwen_doc_tts_models()
            logger.warning("[TTS][QWEN] doc fallback models_count=%s", len(model_ids))
        except Exception as e:
            logger.warning("[TTS][QWEN] doc fallback failed err=%s", str(e))
    logger.warning("[TTS][QWEN] catalog models_count=%s", len(model_ids))

    voices_by_model: dict[str, list[str]] = {}
    for model_id in model_ids[:60]:
        detail_url = root.rstrip("/") + f"/models/{quote(model_id, safe='')}"
        try:
            d_resp = requests.get(detail_url, headers=headers, timeout=20)
            logger.warning("[TTS][QWEN] model detail http_status=%s model=%s", d_resp.status_code, model_id)
            if d_resp.status_code >= 400:
                allowed = _voices_allowed_by_doc(model_id)
                if allowed:
                    fallback_voices = sorted(list(allowed))
                    logger.warning(
                        "[TTS][QWEN] model voices fallback to doc model=%s allowed=%s",
                        model_id,
                        len(fallback_voices),
                    )
                    voices_by_model[model_id] = fallback_voices
                continue
            d_payload = d_resp.json()
            voices = _extract_voices_from_model_meta(d_payload)
            # Strict filter: only voices declared in Bailian docs for this model.
            allowed = _voices_allowed_by_doc(model_id)
            before_filter = len(voices)
            if allowed:
                voices = [v for v in voices if v in allowed]
            if not voices and allowed:
                # Keep strict-doc behavior while avoiding empty lists when model detail lacks voice metadata.
                voices = sorted(list(allowed))
            logger.warning(
                "[TTS][QWEN] model voices parsed model=%s before=%s allowed=%s after=%s",
                model_id,
                before_filter,
                len(allowed),
                len(voices),
            )
            if voices:
                logger.warning(
                    "[TTS][QWEN] model_detail voices model=%s count=%s strict_doc=true",
                    model_id,
                    len(voices),
                )
                voices_by_model[model_id] = voices
        except Exception as e:
            logger.warning("[TTS][QWEN] model detail parse failed model=%s err=%s", model_id, str(e))
            allowed = _voices_allowed_by_doc(model_id)
            if allowed:
                fallback_voices = sorted(list(allowed))
                logger.warning(
                    "[TTS][QWEN] model voices fallback to doc after exception model=%s allowed=%s",
                    model_id,
                    len(fallback_voices),
                )
                voices_by_model[model_id] = fallback_voices
            continue

    return model_ids, voices_by_model


def _normalize_openai_root(base_url: str) -> str:
    b = (base_url or "").strip().rstrip("/")
    if not b:
        return ""
    if b.endswith("/v1"):
        return b
    return b + "/v1"


def _fetch_openai_compatible_tts_models(api_key: str, base_url: str) -> list[str]:
    root = _normalize_openai_root(base_url)
    if not root:
        return []
    url = root + "/models"
    logger.warning("[TTS][OAI] catalog fetch models_url=%s", url)
    resp = requests.get(
        url,
        headers={"Authorization": f"Bearer {api_key}"},
        timeout=20,
    )
    logger.warning("[TTS][OAI] catalog models http_status=%s", resp.status_code)
    if resp.status_code >= 400:
        logger.warning("[TTS][OAI] fetch models failed status=%s body=%s", resp.status_code, resp.text[:200])
        return []
    payload = resp.json()
    model_ids = _extract_tts_model_ids(payload)
    logger.warning("[TTS][OAI] catalog models_count=%s", len(model_ids))
    return model_ids


def _qianwen_tts_to_file(
    api_key: str,
    base_url: str,
    model_name: str,
    text: str,
    voice: str,
    out_file: Path,
) -> None:
    root = _normalize_qianwen_api_root(base_url)
    url = root.rstrip("/") + "/services/aigc/multimodal-generation/generation"
    key_prefix = (api_key[:8] + "...") if api_key else "<empty>"
    logger.warning(
        "[TTS][QWEN] request url=%s model=%s voice=%s text_len=%s key=%s",
        url,
        model_name,
        voice,
        len(text or ""),
        key_prefix,
    )
    allowed_voices = _voices_allowed_by_doc(model_name)
    if allowed_voices and voice not in allowed_voices:
        new_voice = sorted(list(allowed_voices))[0]
        logger.warning(
            "[TTS][QWEN] voice auto-fallback model=%s from=%s to=%s allowed_count=%s",
            model_name,
            voice,
            new_voice,
            len(allowed_voices),
        )
        voice = new_voice
    # 千问 TTS API：输入长度为 UTF-8 字节数 [0, 600]，按字节切分
    _QWEN_TTS_MAX_BYTES = 600

    def _utf8_byte_len(s: str) -> int:
        return len((s or "").encode("utf-8"))

    def _truncate_to_utf8_bytes(s: str, max_bytes: int) -> str:
        b = s.encode("utf-8")
        if len(b) <= max_bytes:
            return s
        return b[:max_bytes].decode("utf-8", errors="ignore").strip() or s[:1]

    def _split_text_for_qwen_tts(full_text: str, max_bytes: int = _QWEN_TTS_MAX_BYTES) -> list[str]:
        txt = (full_text or "").strip()
        if not txt:
            return []
        if _utf8_byte_len(txt) <= max_bytes:
            return [txt]
        parts: list[str] = []
        buf = ""
        for seg in re.split(r"([。！？!?；;，,\n])", txt):
            if not seg:
                continue
            seg_bytes = _utf8_byte_len(seg)
            buf_bytes = _utf8_byte_len(buf)
            if buf_bytes + seg_bytes <= max_bytes:
                buf += seg
                continue
            if buf.strip():
                parts.append(buf.strip())
                buf = ""
            while _utf8_byte_len(seg) > max_bytes:
                parts.append(_truncate_to_utf8_bytes(seg, max_bytes))
                remainder = seg.encode("utf-8")[max_bytes:].decode("utf-8", errors="ignore").strip()
                seg = remainder
            buf = seg
        if buf.strip():
            parts.append(buf.strip())
        return [p for p in parts if p]

    def _qwen_tts_request_once(input_text: str) -> bytes:
        # 发送前按 UTF-8 字节截断，确保符合 API [0, 600] 字节要求
        if _utf8_byte_len(input_text) > _QWEN_TTS_MAX_BYTES:
            input_text = _truncate_to_utf8_bytes(input_text, _QWEN_TTS_MAX_BYTES)
            logger.warning("[TTS][QWEN] request text truncated to 600 utf8 bytes")
        payload = {
            "model": model_name,
            "input": {
                "text": input_text,
                "voice": voice,
                "language_type": "Chinese",
            },
        }
        resp = requests.post(
            url,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=payload,
            timeout=120,
        )
        if resp.status_code >= 400:
            request_id = ""
            code = ""
            message = resp.text[:600]
            try:
                err_data = resp.json()
                request_id = str(err_data.get("request_id") or "")
                code = str(err_data.get("code") or "")
                message = str(err_data.get("message") or message)
            except Exception:
                pass
            logger.error(
                "[TTS][QWEN] failed status=%s code=%s request_id=%s url=%s model=%s text_len=%s message=%s",
                resp.status_code,
                code,
                request_id,
                url,
                model_name,
                len(input_text),
                message,
            )
            raise RuntimeError(
                f"HTTP {resp.status_code} code={code} request_id={request_id} model={model_name} url={url} message={message}"
            )

        data = resp.json()
        audio = (data.get("output") or {}).get("audio") or {}
        audio_url = audio.get("url")
        audio_b64 = audio.get("data")
        if audio_url:
            logger.warning("[TTS][QWEN] audio_url returned model=%s url=%s", model_name, audio_url)
            raw = requests.get(audio_url, timeout=120)
            if raw.status_code >= 400:
                raise RuntimeError(f"下载音频失败 HTTP {raw.status_code}: {raw.text[:300]}")
            return raw.content
        if audio_b64:
            logger.warning("[TTS][QWEN] audio_base64 returned model=%s bytes_b64=%s", model_name, len(audio_b64))
            return base64.b64decode(audio_b64)
        logger.error(
            "[TTS][QWEN] output missing audio field model=%s raw=%s",
            model_name,
            json.dumps(data, ensure_ascii=False)[:1200],
        )
        raise RuntimeError(f"未拿到音频输出: {json.dumps(data, ensure_ascii=False)[:800]}")

    chunks = _split_text_for_qwen_tts(text, max_bytes=_QWEN_TTS_MAX_BYTES)
    if not chunks:
        raise RuntimeError("文本为空，无法进行 TTS")
    logger.warning("[TTS][QWEN] split into chunks count=%s lens=%s", len(chunks), [len(x) for x in chunks[:20]])
    if len(chunks) == 1:
        out_file.write_bytes(_qwen_tts_request_once(chunks[0]))
        return

    # Multi-chunk synthesis: render each part and merge with ffmpeg.
    with tempfile.TemporaryDirectory(prefix="qwen_tts_") as td:
        tmp_dir = Path(td)
        part_files: list[Path] = []
        for i, chunk in enumerate(chunks, start=1):
            part = tmp_dir / f"part_{i:03d}.audio"
            part.write_bytes(_qwen_tts_request_once(chunk))
            part_files.append(part)
        list_file = tmp_dir / "concat.txt"
        lines: list[str] = []
        for p in part_files:
            escaped = str(p).replace("'", "'\\''")
            lines.append(f"file '{escaped}'")
        list_file.write_text("\n".join(lines), encoding="utf-8")
        cmd = [
            "ffmpeg",
            "-y",
            "-f",
            "concat",
            "-safe",
            "0",
            "-i",
            str(list_file),
            "-vn",
            "-acodec",
            "libmp3lame",
            str(out_file),
        ]
        logger.warning("[TTS][QWEN] merge parts via ffmpeg parts=%s out=%s", len(part_files), out_file)
        proc = subprocess.run(cmd, capture_output=True, text=True)
        if proc.returncode != 0:
            logger.error("[TTS][QWEN] ffmpeg merge failed rc=%s stderr=%s", proc.returncode, (proc.stderr or "")[:1200])
            raise RuntimeError(f"音频合并失败（ffmpeg rc={proc.returncode}）: {(proc.stderr or '')[:500]}")


def _llm_generate_slides(text: str, max_slides: int = 20) -> list[dict[str, Any]]:
    from ..rag.config import get_rag_settings
    from ..rag.llm import get_llm

    system = (
        "你是教学PPT助手。请根据教材文本生成教学幻灯片内容。"
        "输出 JSON 数组，每项包含 slide_no(int), title(str), bullets(str[]), notes(str)。"
        f"总页数不超过 {max_slides}。bullets 每页 3~5 条，notes 为讲解备注。"
    )
    prompt = f"{system}\n\n教材文本：\n{text[:12000]}"
    settings = get_rag_settings()
    raw = get_llm(settings).generate(
        prompt,
        max_tokens=max(settings.llm_max_tokens, 4096),
        temperature=settings.llm_temperature if settings.llm_temperature is not None else 0.3,
    )
    raw = (raw or "").strip()
    m = re.search(r"\[\s*\{[\s\S]*\}\s*\]", raw)
    payload = m.group(0) if m else raw
    data = json.loads(payload)
    if not isinstance(data, list):
        raise RuntimeError("LLM 输出格式不正确")
    slides: list[dict[str, Any]] = []
    for i, item in enumerate(data[:max_slides], start=1):
        title = str(item.get("title") or "").strip()[:40] or f"教学要点{i}"
        bullets = item.get("bullets") if isinstance(item.get("bullets"), list) else []
        bullets_clean = [str(x).strip()[:80] for x in bullets if str(x).strip()][:5]
        notes = str(item.get("notes") or "").strip()[:500]
        slides.append(
            {"slide_no": i, "title": title, "bullets": bullets_clean, "notes": notes}
        )
    if not slides:
        raise RuntimeError("LLM 未返回可用页内容")
    return slides


def _create_ppt_from_slides(
    slides: list[dict[str, Any]],
    output_path: Path,
    template_path: Path | None = None,
) -> int:
    prs = Presentation(str(template_path)) if template_path and template_path.exists() else Presentation()
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    for s in slides:
        slide = prs.slides.add_slide(layout)
        title_box = slide.shapes.title
        if title_box:
            title_box.text = s.get("title", "")
        body_ph = None
        for ph in slide.placeholders:
            if getattr(ph, "is_placeholder", False) and hasattr(ph, "text_frame"):
                if slide.shapes.title is not None and ph == slide.shapes.title:
                    continue
                body_ph = ph
                break
        if body_ph and hasattr(body_ph, "text_frame"):
            tf = body_ph.text_frame
            tf.clear()
            bullets = s.get("bullets", [])
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(b)
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(24)
        notes = (s.get("notes") or "").strip()
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return len(prs.slides)


def _extract_ppt_notes_and_text(ppt_path: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(ppt_path))
    rows: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
        shape_texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                cleaned = shape.text.strip()
                if cleaned:
                    shape_texts.append(cleaned)
        rows.append(
            {
                "slide_no": idx,
                "title": shape_texts[0] if shape_texts else f"第{idx}页",
                "content_text": "\n".join(shape_texts),
                "notes_text": notes_text,
            }
        )
    return rows


def _fallback_script(segments: list[dict[str, Any]]) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for seg in segments:
        slide_no = seg["slide_no"]
        base = (seg.get("notes_text") or "").strip() or (seg.get("content_text") or "").strip()
        base = re.sub(r"\s+", " ", base)
        if len(base) > 320:
            base = base[:320] + "。"
        text = f"第{slide_no}页。{base}" if base else f"第{slide_no}页，我们讲解本页核心概念。"
        out.append({"slide_no": slide_no, "script": text})
    return out


def _llm_script(segments: list[dict[str, Any]]) -> list[dict[str, Any]]:
    from ..rag.config import get_rag_settings
    from ..rag.llm import get_llm

    system = (
        "你是教师讲解稿助手。请按页生成讲解词。"
        "输出 JSON 数组，每项包含 slide_no(int), script(str)。语气自然、简洁、适合课堂讲解。"
    )
    src = json.dumps(segments, ensure_ascii=False)[:12000]
    prompt = f"{system}\n\nPPT 按页内容：\n{src}"
    settings = get_rag_settings()
    raw = get_llm(settings).generate(
        prompt,
        max_tokens=max(settings.llm_max_tokens, 4096),
        temperature=0.4,
    )
    raw = (raw or "").strip()
    m = re.search(r"\[\s*\{[\s\S]*\}\s*\]", raw)
    payload = m.group(0) if m else raw
    data = json.loads(payload)
    if not isinstance(data, list):
        raise RuntimeError("LLM 输出格式不正确")
    out: list[dict[str, Any]] = []
    for item in data:
        slide_no = int(item.get("slide_no") or 0)
        script = str(item.get("script") or "").strip()
        if slide_no > 0 and script:
            out.append({"slide_no": slide_no, "script": script})
    if not out:
        raise RuntimeError("LLM 未返回可用讲解稿")
    return out


def _require_tool(cmd: str, hint: str) -> None:
    if shutil.which(cmd):
        return
    raise HTTPException(status_code=400, detail=f"缺少依赖命令 `{cmd}`，请先安装后重试（{hint}）")


def _run_cmd(args: list[str], cwd: Path | None = None) -> None:
    try:
        subprocess.run(args, cwd=str(cwd) if cwd else None, check=True, capture_output=True, text=True)
    except subprocess.CalledProcessError as e:
        stderr = (e.stderr or "").strip()[:500]
        raise HTTPException(status_code=400, detail=f"命令执行失败: {' '.join(args)}; {stderr}")


def _get_audio_duration_seconds(audio_path: Path) -> float | None:
    """用 ffprobe 获取音频时长（秒），失败返回 None。"""
    try:
        out = subprocess.run(
            [
                "ffprobe",
                "-v",
                "error",
                "-show_entries",
                "format=duration",
                "-of",
                "default=noprint_wrappers=1:nokey=1",
                str(audio_path),
            ],
            capture_output=True,
            text=True,
            timeout=10,
        )
        if out.returncode == 0 and out.stdout:
            return float(out.stdout.strip())
    except (subprocess.TimeoutExpired, ValueError, FileNotFoundError):
        pass
    return None


def _build_stage6_video(
    ppt_path: Path,
    audio_path: Path,
    out_video: Path,
    durations: list[float] | None = None,
    script_segments_path: Path | None = None,
) -> dict[str, Any]:
    _require_tool("soffice", "用于将 PPTX 转 PDF")
    _require_tool("pdftoppm", "用于将 PDF 转图片")
    _require_tool("ffmpeg", "用于合成视频")

    with tempfile.TemporaryDirectory(prefix="pipeline_stage6_") as tmpdir:
        tmp = Path(tmpdir)
        # 1) PPTX -> PDF
        _run_cmd(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(tmp),
                str(ppt_path),
            ]
        )
        pdf_path = tmp / f"{ppt_path.stem}.pdf"
        if not pdf_path.exists():
            raise HTTPException(status_code=400, detail="PPT 转 PDF 失败，未生成 PDF")

        # 2) PDF -> PNG sequence
        img_prefix = tmp / "slide"
        _run_cmd(["pdftoppm", "-png", str(pdf_path), str(img_prefix)])
        images = sorted(tmp.glob("slide-*.png"))
        if not images:
            raise HTTPException(status_code=400, detail="PDF 转图片失败，未生成页图")

        # 3) 每页时长：有传入则用传入的；否则按「每页字数/总字数 * 音频总时长」分配，无分段文件则均分
        default_per_slide = 8.0
        if durations is not None and len(durations) > 0:
            durs = list(durations)
        else:
            total_sec = _get_audio_duration_seconds(audio_path)
            if total_sec is None or total_sec <= 0:
                durs = [default_per_slide] * len(images)
            else:
                segments: list[dict[str, Any]] = []
                if script_segments_path and script_segments_path.exists():
                    try:
                        raw = json.loads(script_segments_path.read_text(encoding="utf-8"))
                        if isinstance(raw, list):
                            segments = [x for x in raw if isinstance(x, dict)]
                    except (json.JSONDecodeError, OSError):
                        pass
                n_slides = len(images)
                use_segments = segments and len(segments) > 0
                if use_segments:
                    # 取前 n_slides 段，按字数占比分配时长
                    segs = segments[:n_slides]
                    total_chars = sum(len(str(s.get("script", "") or "")) for s in segs)
                    if total_chars > 0:
                        durs = [
                            max(1.0, len(str(s.get("script", "") or "")) / total_chars * total_sec)
                            for s in segs
                        ]
                        if len(durs) < n_slides:
                            assigned = sum(durs)
                            remaining = max(0, total_sec - assigned)
                            pad = remaining / (n_slides - len(durs)) if (n_slides - len(durs)) else default_per_slide
                            durs = durs + [max(1.0, pad)] * (n_slides - len(durs))
                        else:
                            durs = durs[:n_slides]
                    else:
                        use_segments = False
                if not use_segments:
                    per_slide = total_sec / n_slides
                    durs = [max(1.0, per_slide)] * n_slides
        if len(durs) < len(images):
            durs = durs + [durs[-1] if durs else default_per_slide] * (len(images) - len(durs))
        durs = durs[: len(images)]

        seg_dir = tmp / "segments"
        seg_dir.mkdir(parents=True, exist_ok=True)
        seg_paths: list[Path] = []
        for i, img in enumerate(images, start=1):
            seg = seg_dir / f"seg_{i:03d}.mp4"
            _run_cmd(
                [
                    "ffmpeg",
                    "-y",
                    "-loop",
                    "1",
                    "-i",
                    str(img),
                    "-t",
                    f"{max(1.0, float(durs[i - 1])):.2f}",
                    "-vf",
                    "scale=1280:720:force_original_aspect_ratio=decrease,pad=1280:720:(ow-iw)/2:(oh-ih)/2",
                    "-r",
                    "30",
                    "-pix_fmt",
                    "yuv420p",
                    str(seg),
                ]
            )
            seg_paths.append(seg)

        concat_txt = tmp / "segments.txt"
        concat_txt.write_text(
            "\n".join([f"file '{p.as_posix()}'" for p in seg_paths]) + "\n",
            encoding="utf-8",
        )
        merged_video = tmp / "slides.mp4"
        _run_cmd(
            [
                "ffmpeg",
                "-y",
                "-f",
                "concat",
                "-safe",
                "0",
                "-i",
                str(concat_txt),
                "-c",
                "copy",
                str(merged_video),
            ]
        )

        out_video.parent.mkdir(parents=True, exist_ok=True)
        _run_cmd(
            [
                "ffmpeg",
                "-y",
                "-i",
                str(merged_video),
                "-i",
                str(audio_path),
                "-c:v",
                "libx264",
                "-c:a",
                "aac",
                "-shortest",
                str(out_video),
            ]
        )
        return {
            "slide_count": len(images),
            "durations": durs,
        }


class PipelineCreateIn(BaseModel):
    title: str = Field(default="新教学流水线", min_length=1, max_length=80)
    chapter_id: int | None = None


class PipelineOut(BaseModel):
    workflow_id: str
    title: str
    chapter_id: int | None
    created_at: str
    updated_at: str
    stages: dict[str, Any]


class Stage2RunIn(BaseModel):
    source_file: str = "stage1/extracted_content.md"
    max_slides: int = Field(default=20, ge=5, le=60)
    prefer_llm: bool = True
    title: str = "课程教学PPT"
    output_json_file: str = "stage2/slides_content.json"


class Stage3RunIn(BaseModel):
    source_file: str = "stage2/slides_content.json"
    output_file: str = "stage3/generated_draft.pptx"
    template_file: str | None = None


class Stage4RunIn(BaseModel):
    source_ppt: str = "stage3/edited.pptx"
    fallback_ppt: str = "stage3/generated_draft.pptx"
    prefer_llm: bool = True
    output_segments_file: str = "stage4/script_segments.json"
    output_script_file: str = "stage4/narration_script.txt"


class Stage5RunIn(BaseModel):
    script_file: str = "stage4/narration_script.txt"
    output_file: str = "stage5/full_narration.mp3"
    model: str = ""  # 空则使用 RAG 配置的默认 TTS（admin/rag 默认 TTS）
    voice: str = "alloy"
    speed: float = Field(default=1.0, ge=0.5, le=1.5)


class TTSPreviewIn(BaseModel):
    """试听请求：使用 RAG 默认 TTS，生成固定试听文案的短音频。"""
    voice: str = "alloy"
    speed: float = Field(default=1.0, ge=0.5, le=1.5)


class Stage6RunIn(BaseModel):
    ppt_file: str = "stage3/edited.pptx"
    fallback_ppt: str = "stage3/generated_draft.pptx"
    audio_file: str = "stage5/full_narration.mp3"
    output_file: str = "stage6/final_video.mp4"
    timing_file: str | None = "stage6/timeline.json"
    script_segments_file: str | None = "stage4/script_segments.json"  # 用于按每页字数占比分配时长
    default_slide_seconds: float = Field(default=8.0, ge=1.0, le=60.0)


class TeacherPdfDocOut(BaseModel):
    id: int
    chapter_id: int | None
    chapter_title: str | None
    course_id: int | None
    course_name: str | None
    title: str
    file_name: str | None
    file_size: int | None = None
    parse_status: str | None
    created_at: str | None
    workflow_id: str


class FileTextOut(BaseModel):
    path: str
    content: str


class FileTextIn(BaseModel):
    path: str
    content: str


class TTSModelOptionOut(BaseModel):
    value: str
    label: str
    provider_id: str
    provider_name: str
    provider_type: str
    model: str


class TTSVoiceOptionOut(BaseModel):
    value: str
    label: str
    gender: str | None = None  # "male" | "female"，未知则为 None


class TTSModelsOut(BaseModel):
    default_model: str
    options: list[TTSModelOptionOut]
    voices_by_provider_type: dict[str, list[TTSVoiceOptionOut]]
    voices_by_model: dict[str, list[TTSVoiceOptionOut]]


@router.get("/pdf-docs", response_model=list[TeacherPdfDocOut])
async def list_teacher_pdf_docs(
    _q: str | None = Query(None, alias="q"),
    user: User = Depends(require_teacher),
):
    rows = _read_pdf_docs_index()
    out: list[TeacherPdfDocOut] = []
    for row in sorted(rows, key=lambda x: int(x.get("id") or 0), reverse=True):
        if int(row.get("owner_id") or 0) != user.id:
            continue
        out.append(
            TeacherPdfDocOut(
                id=int(row.get("id") or 0),
                chapter_id=None,
                chapter_title=None,
                course_id=None,
                course_name=None,
                title=str(row.get("title") or ""),
                file_name=row.get("file_name"),
                file_size=int(row.get("file_size") or 0) or None,
                parse_status=str(row.get("parse_status") or "uploaded"),
                created_at=str(row.get("created_at") or ""),
                workflow_id=_workflow_id_for_doc(int(row.get("id") or 0)),
            )
        )
    return out


@router.post("/pdf-docs/upload", response_model=TeacherPdfDocOut)
async def upload_pipeline_pdf_doc(
    file: UploadFile = File(...),
    user: User = Depends(require_teacher),
):
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="请上传 PDF 文件")
    data = await file.read()
    if not data:
        raise HTTPException(status_code=400, detail="文件内容为空")

    rows = _read_pdf_docs_index()
    new_id = (max([int(x.get("id") or 0) for x in rows]) + 1) if rows else 1
    safe_name = _safe_name(file.filename)
    saved_name = f"{new_id}_{int(datetime.utcnow().timestamp())}_{safe_name}"
    rel_path = f"docs/{saved_name}"
    abs_path = _pipeline_root() / rel_path
    abs_path.parent.mkdir(parents=True, exist_ok=True)
    abs_path.write_bytes(data)

    row = {
        "id": new_id,
        "owner_id": user.id,
        "title": file.filename,
        "file_name": file.filename,
        "file_path": rel_path,
        "file_size": len(data),
        "parse_status": "uploaded",
        "created_at": _now_iso(),
    }
    rows.append(row)
    _write_pdf_docs_index(rows)
    return TeacherPdfDocOut(
        id=new_id,
        chapter_id=None,
        chapter_title=None,
        course_id=None,
        course_name=None,
        title=file.filename,
        file_name=file.filename,
        file_size=len(data),
        parse_status="uploaded",
        created_at=row["created_at"],
        workflow_id=_workflow_id_for_doc(new_id),
    )


@router.get("/tts-models", response_model=TTSModelsOut)
async def list_tts_models(
    user: User = Depends(require_teacher),
):
    _ = user
    try:
        from ..rag.config_store import get_providers_list_raw, get_default_llm, get_default_tts
        providers = get_providers_list_raw()
        default_tts = (get_default_tts() or "").strip()
        default_llm = (get_default_llm() or "").strip()
    except Exception:
        providers = []
        default_tts = ""
        default_llm = ""

    logger.warning("[TTS] list_tts_models providers_count=%s default_tts=%s default_llm=%s", len(providers), default_tts, default_llm)
    options: list[TTSModelOptionOut] = []
    voices_by_provider_type: dict[str, list[TTSVoiceOptionOut]] = {}
    voices_by_model: dict[str, list[TTSVoiceOptionOut]] = {}
    for p in providers:
        provider_id = (p.get("id") or "").strip()
        provider_name = (p.get("name") or "").strip() or provider_id
        provider_type = (p.get("type") or "").strip() or "openai_compatible"
        api_key = (p.get("api_key") or "").strip()
        base_url = (p.get("base_url") or "").strip()
        logger.warning(
            "[TTS] provider start id=%s type=%s name=%s has_key=%s base_url=%s",
            provider_id,
            provider_type,
            provider_name,
            bool(api_key),
            base_url or "<empty>",
        )
        if not provider_id:
            continue
        models: list[str] = []
        model_voice_map: dict[str, list[str]] = {}
        try:
            if api_key:
                if provider_type == "qianwen":
                    models, model_voice_map = _fetch_qianwen_tts_catalog(api_key=api_key, base_url=base_url)
                else:
                    models = _fetch_openai_compatible_tts_models(api_key=api_key, base_url=base_url)
            else:
                logger.warning("[TTS] provider skipped (no api key) id=%s type=%s", provider_id, provider_type)
        except Exception as e:
            logger.warning(
                "[TTS] fetch catalog failed provider=%s type=%s err=%s",
                provider_id,
                provider_type,
                str(e),
            )
        logger.warning(
            "[TTS] provider result id=%s type=%s models=%s models_with_voices=%s",
            provider_id,
            provider_type,
            len(models),
            len(model_voice_map),
        )

        for m in models:
            model_key = f"{provider_id}:{m}"
            options.append(
                TTSModelOptionOut(
                    value=model_key,
                    label=f"{provider_name} - {m}",
                    provider_id=provider_id,
                    provider_name=provider_name,
                    provider_type=provider_type,
                    model=m,
                )
            )

            model_voices = model_voice_map.get(m, [])
            if model_voices:
                voices_by_model[model_key] = [
                    TTSVoiceOptionOut(value=v, label=_voice_label_zh(v), gender=_voice_gender(v))
                    for v in model_voices
                ]
                existing = {x.value for x in voices_by_provider_type.get(provider_type, [])}
                merged = list(voices_by_provider_type.get(provider_type, []))
                for v in model_voices:
                    if v in existing:
                        continue
                    merged.append(TTSVoiceOptionOut(value=v, label=_voice_label_zh(v), gender=_voice_gender(v)))
                    existing.add(v)
                voices_by_provider_type[provider_type] = merged

    default_model = (
        default_tts
        if ":" in default_tts
        else default_llm
        if ":" in default_llm
        else (options[0].value if options else "gpt-4o-mini-tts")
    )
    logger.warning(
        "[TTS] list_tts_models done options=%s voices_by_provider_type=%s voices_by_model=%s default_model=%s",
        len(options),
        sum(len(v) for v in voices_by_provider_type.values()),
        len(voices_by_model),
        default_model,
    )
    return TTSModelsOut(
        default_model=default_model,
        options=options,
        voices_by_provider_type=voices_by_provider_type,
        voices_by_model=voices_by_model,
    )


# 试听固定文案，用于前端「试听」按钮
TTS_PREVIEW_TEXT = "你好，这是当前音色的试听。"


def _tts_preview_cache_dir() -> Path:
    d = _pipeline_root() / "tts_preview_cache"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _tts_preview_cache_path(model_in: str, voice: str, speed: float) -> Path:
    key = f"{model_in}|{voice}|{speed}"
    name = hashlib.sha256(key.encode("utf-8")).hexdigest()[:32] + ".mp3"
    return _tts_preview_cache_dir() / name


@router.post("/tts-preview")
async def tts_preview(
    body: TTSPreviewIn,
    user: User = Depends(require_teacher),
) -> Response:
    """使用 RAG 默认 TTS 生成试听短音频，返回 audio/mpeg。相同 model+voice+speed 会走缓存，不重复调用模型。"""
    _ = user
    model_in = ""
    try:
        from ..rag.config_store import get_default_tts
        model_in = (get_default_tts() or "").strip()
    except Exception:
        pass
    if not model_in:
        model_in = "gpt-4o-mini-tts"

    cache_path = _tts_preview_cache_path(model_in, body.voice, body.speed)
    if cache_path.exists():
        return Response(content=cache_path.read_bytes(), media_type="audio/mpeg")

    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
        out = Path(tmp.name)
    try:
        provider, resolved_model_name = _resolve_provider_and_model(model_in)
        provider_type = ((provider or {}).get("type") or "").strip()
        if provider and provider_type == "qianwen":
            api_key = ((provider or {}).get("api_key") or "").strip()
            if not api_key:
                raise HTTPException(status_code=400, detail="该模型提供商未配置 API Key")
            base_url = ((provider or {}).get("base_url") or "").strip()
            _qianwen_tts_to_file(
                api_key=api_key,
                base_url=base_url,
                model_name=resolved_model_name,
                text=TTS_PREVIEW_TEXT,
                voice=body.voice,
                out_file=out,
            )
        else:
            client, model_name = _resolve_tts_client_and_model(model_in)
            speech = client.audio.speech.create(
                model=model_name,
                voice=body.voice,
                input=TTS_PREVIEW_TEXT,
                speed=body.speed,
            )
            speech.stream_to_file(str(out))
        audio_bytes = out.read_bytes()
        cache_path.write_bytes(audio_bytes)
        return Response(content=audio_bytes, media_type="audio/mpeg")
    except Exception as e:
        logger.exception("[TTS] preview failed voice=%s", body.voice)
        raise HTTPException(status_code=400, detail=f"试听生成失败: {str(e)}") from e
    finally:
        if out.exists():
            try:
                out.unlink()
            except Exception:
                pass


@router.post("/pdf-docs/{doc_id}/workflow", response_model=PipelineOut)
async def get_or_create_pdf_workflow(
    doc_id: int,
    user: User = Depends(require_teacher),
):
    doc = _find_pdf_doc_or_404(user.id, doc_id)
    rel_path = str(doc.get("file_path") or "")
    if not rel_path:
        raise HTTPException(status_code=400, detail="PDF 原文件路径为空")
    abs_pdf = _pipeline_root() / rel_path
    if not abs_pdf.exists():
        raise HTTPException(status_code=400, detail="PDF 原文件不存在")

    workflow_id = _workflow_id_for_doc(doc_id)
    wf_dir = _ensure_workflow_dirs(workflow_id)
    manifest_file = wf_dir / "manifest.json"
    if not manifest_file.exists():
        manifest = {
            "workflow_id": workflow_id,
            "title": str(doc.get("file_name") or doc.get("title") or f"PDF-{doc_id}"),
            "chapter_id": None,
            "source_doc_id": doc_id,
            "source_doc_name": str(doc.get("file_name") or doc.get("title") or ""),
            "source_doc_path": rel_path,
            "created_at": _now_iso(),
            "updated_at": _now_iso(),
            "owner_hint": "teacher",
            "stages": {
                f"stage{i}": {"status": "pending", "updated_at": _now_iso(), "outputs": []}
                for i in range(1, 7)
            },
        }
        manifest_file.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    # 同步源 PDF 到阶段1输入
    src_pdf = wf_dir / "stage1" / "source.pdf"
    if not src_pdf.exists():
        src_pdf.write_bytes(abs_pdf.read_bytes())

    return PipelineOut(**_read_manifest(workflow_id))


@router.post("/workflows", response_model=PipelineOut)
async def create_workflow(
    body: PipelineCreateIn,
    user: User = Depends(require_teacher),
):
    _ = user
    workflow_id = uuid.uuid4().hex[:16]
    wf_dir = _pipeline_root() / workflow_id
    wf_dir.mkdir(parents=True, exist_ok=True)
    for sub in ["stage1", "stage2", "stage3", "stage4", "stage5", "stage6"]:
        (wf_dir / sub).mkdir(parents=True, exist_ok=True)

    manifest = {
        "workflow_id": workflow_id,
        "title": body.title.strip(),
        "chapter_id": body.chapter_id,
        "created_at": _now_iso(),
        "updated_at": _now_iso(),
        "owner_hint": "teacher",
        "stages": {
            f"stage{i}": {"status": "pending", "updated_at": _now_iso(), "outputs": []}
            for i in range(1, 7)
        },
    }
    (wf_dir / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return PipelineOut(**manifest)


@router.get("/workflows/{workflow_id}", response_model=PipelineOut)
async def get_workflow(
    workflow_id: str,
    user: User = Depends(require_teacher),
):
    _ = user
    manifest = _read_manifest(workflow_id)
    return PipelineOut(**manifest)


@router.post("/workflows/{workflow_id}/stage1/extract")
async def stage1_extract(
    workflow_id: str,
    file: UploadFile | None = File(None),
    user: User = Depends(require_teacher),
):
    _ = user
    manifest = _read_manifest(workflow_id)
    wf_dir = _workflow_dir(workflow_id)
    stage_dir = wf_dir / "stage1"
    stage_dir.mkdir(parents=True, exist_ok=True)
    src_path = stage_dir / "source.pdf"

    pdf_bytes: bytes
    if file is not None:
        if not file.filename or not file.filename.lower().endswith(".pdf"):
            raise HTTPException(status_code=400, detail="请上传 PDF 文件")
        pdf_bytes = await file.read()
        src_path.write_bytes(pdf_bytes)
    else:
        source_doc_path = str(manifest.get("source_doc_path") or "").strip()
        source_doc_id = int(manifest.get("source_doc_id") or 0)
        if not source_doc_path and not source_doc_id:
            raise HTTPException(status_code=400, detail="该工作流未绑定 PDF，请上传 PDF")
        if not source_doc_path and source_doc_id:
            doc = _find_pdf_doc_or_404(user.id, source_doc_id)
            source_doc_path = str(doc.get("file_path") or "")
        abs_pdf = _pipeline_root() / source_doc_path
        if not abs_pdf.exists():
            raise HTTPException(status_code=400, detail="PDF 原文件不存在")
        pdf_bytes = abs_pdf.read_bytes()
        src_path.write_bytes(pdf_bytes)

    text, total_pages = _extract_pdf_text(pdf_bytes)
    if not text.strip():
        raise HTTPException(status_code=400, detail="未提取到可用文本，请先人工 OCR 后重试")

    extracted_file = stage_dir / "extracted_content.md"
    extracted_file.write_text(text, encoding="utf-8")

    chapter_outputs: list[str] = []
    chapter_meta: list[dict[str, Any]] = []
    chapter_dir = stage_dir / "chapters"
    chapter_dir.mkdir(parents=True, exist_ok=True)
    chapter_splits = _split_text_by_chapter(text)
    for ch in chapter_splits:
        no = int(ch["chapter_no"])
        title = str(ch["title"])
        content = str(ch["content"]).strip()
        rel = f"stage1/chapters/chapter_{no:02d}.md"
        fp = _safe_workflow_file(workflow_id, rel)
        fp.write_text(content + "\n", encoding="utf-8")
        chapter_outputs.append(rel)
        chapter_meta.append(
            {
                "chapter_no": no,
                "title": title,
                "path": rel,
                "char_count": len(content),
            }
        )
    chapter_manifest_rel = "stage1/chapter_manifest.json"
    chapter_manifest_file = _safe_workflow_file(workflow_id, chapter_manifest_rel)
    chapter_manifest_file.write_text(
        json.dumps({"chapters": chapter_meta}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    chapter_outputs.append(chapter_manifest_rel)

    _touch_stage(
        workflow_id,
        "stage1",
        "done",
        outputs=["stage1/source.pdf", "stage1/extracted_content.md", *chapter_outputs],
        extra={"page_count": total_pages, "chapter_count": len(chapter_meta)},
    )
    return {
        "ok": True,
        "workflow_id": workflow_id,
        "page_count": total_pages,
        "chapter_count": len(chapter_meta),
        "chapters": chapter_meta,
        "outputs": ["stage1/source.pdf", "stage1/extracted_content.md", *chapter_outputs],
    }


@router.post("/workflows/{workflow_id}/stage2/generate")
async def stage2_generate(
    workflow_id: str,
    body: Stage2RunIn,
    user: User = Depends(require_teacher),
):
    _ = user
    src = _safe_workflow_file(workflow_id, body.source_file)
    if not src.exists():
        raise HTTPException(status_code=404, detail=f"文件不存在: {body.source_file}")
    text = src.read_text(encoding="utf-8").strip()
    if not text:
        raise HTTPException(status_code=400, detail="输入文本为空")

    err: str | None = None
    slides: list[dict[str, Any]]
    if body.prefer_llm:
        try:
            slides = _llm_generate_slides(text, max_slides=body.max_slides)
        except Exception as e:
            err = str(e)
            slides = _fallback_build_slides(text, max_slides=body.max_slides)
    else:
        slides = _fallback_build_slides(text, max_slides=body.max_slides)

    wf_dir = _workflow_dir(workflow_id)
    slides_json = {"title": body.title.strip(), "slides": slides}
    json_file = _safe_workflow_file(workflow_id, body.output_json_file)
    json_file.parent.mkdir(parents=True, exist_ok=True)
    json_file.write_text(json.dumps(slides_json, ensure_ascii=False, indent=2), encoding="utf-8")
    _touch_stage(
        workflow_id,
        "stage2",
        "done",
        outputs=[body.output_json_file],
        extra={"slide_count": len(slides)},
    )
    return {
        "ok": True,
        "workflow_id": workflow_id,
        "slide_count": len(slides),
        "fallback_used": bool(err),
        "fallback_reason": err,
        "outputs": [body.output_json_file],
    }


@router.post("/workflows/{workflow_id}/stage3/generate-ppt")
async def stage3_generate_ppt(
    workflow_id: str,
    body: Stage3RunIn,
    user: User = Depends(require_teacher),
):
    _ = user
    src = _safe_workflow_file(workflow_id, body.source_file)
    if not src.exists():
        raise HTTPException(status_code=404, detail=f"文件不存在: {body.source_file}")
    data = json.loads(src.read_text(encoding="utf-8"))
    slides = data.get("slides") if isinstance(data, dict) else None
    if not isinstance(slides, list) or not slides:
        raise HTTPException(status_code=400, detail="slides_content.json 格式不正确")

    output = _safe_workflow_file(workflow_id, body.output_file)
    template = _safe_workflow_file(workflow_id, body.template_file) if body.template_file else None
    count = _create_ppt_from_slides(slides, output_path=output, template_path=template)

    _touch_stage(
        workflow_id,
        "stage3",
        "done",
        outputs=[body.output_file],
        extra={"slide_count": count},
    )
    return {"ok": True, "workflow_id": workflow_id, "slide_count": count, "output": body.output_file}


@router.post("/workflows/{workflow_id}/stage3/upload-edited-ppt")
async def stage3_upload_edited_ppt(
    workflow_id: str,
    file: UploadFile = File(...),
    user: User = Depends(require_teacher),
):
    _ = user
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="请上传 .pptx 文件")
    dst = _safe_workflow_file(workflow_id, "stage3/edited.pptx")
    dst.parent.mkdir(parents=True, exist_ok=True)
    dst.write_bytes(await file.read())
    _touch_stage(workflow_id, "stage3", "done", outputs=["stage3/edited.pptx"])
    return {"ok": True, "workflow_id": workflow_id, "output": "stage3/edited.pptx"}


@router.post("/workflows/{workflow_id}/stage4/generate-script")
async def stage4_generate_script(
    workflow_id: str,
    body: Stage4RunIn,
    user: User = Depends(require_teacher),
):
    _ = user
    ppt = _safe_workflow_file(workflow_id, body.source_ppt)
    if not ppt.exists():
        ppt = _safe_workflow_file(workflow_id, body.fallback_ppt)
    if not ppt.exists():
        raise HTTPException(status_code=404, detail="未找到可用 PPT，请先执行阶段3")

    segments = _extract_ppt_notes_and_text(ppt)
    if not segments:
        raise HTTPException(status_code=400, detail="PPT 没有可读取内容")

    err: str | None = None
    scripts: list[dict[str, Any]]
    if body.prefer_llm:
        try:
            scripts = _llm_script(segments)
        except Exception as e:
            err = str(e)
            scripts = _fallback_script(segments)
    else:
        scripts = _fallback_script(segments)

    seg_file = _safe_workflow_file(workflow_id, body.output_segments_file)
    seg_file.parent.mkdir(parents=True, exist_ok=True)
    seg_file.write_text(json.dumps(scripts, ensure_ascii=False, indent=2), encoding="utf-8")
    merged = "\n\n".join([f"[第{x['slide_no']}页]\n{x['script']}" for x in scripts]) + "\n"
    script_file = _safe_workflow_file(workflow_id, body.output_script_file)
    script_file.parent.mkdir(parents=True, exist_ok=True)
    script_file.write_text(merged, encoding="utf-8")

    _touch_stage(
        workflow_id,
        "stage4",
        "done",
        outputs=[body.output_segments_file, body.output_script_file],
        extra={"segment_count": len(scripts)},
    )
    return {
        "ok": True,
        "workflow_id": workflow_id,
        "segment_count": len(scripts),
        "fallback_used": bool(err),
        "fallback_reason": err,
        "outputs": [body.output_segments_file, body.output_script_file],
    }


@router.post("/workflows/{workflow_id}/stage5/tts")
async def stage5_tts(
    workflow_id: str,
    body: Stage5RunIn,
    user: User = Depends(require_teacher),
):
    _ = user
    src = _safe_workflow_file(workflow_id, body.script_file)
    if not src.exists():
        raise HTTPException(status_code=404, detail=f"文件不存在: {body.script_file}")
    script = src.read_text(encoding="utf-8").strip()
    if not script:
        raise HTTPException(status_code=400, detail="讲解稿为空")

    out = _safe_workflow_file(workflow_id, body.output_file)
    out.parent.mkdir(parents=True, exist_ok=True)
    model_in = (body.model or "").strip()
    if not model_in:
        try:
            from ..rag.config_store import get_default_tts
            model_in = (get_default_tts() or "").strip()
        except Exception:
            model_in = ""
        if not model_in:
            model_in = "gpt-4o-mini-tts"
    try:
        provider, resolved_model_name = _resolve_provider_and_model(model_in)
        provider_type = ((provider or {}).get("type") or "").strip()
        logger.warning(
            "[TTS] stage5 start workflow=%s model_in=%s resolved_model=%s provider_type=%s voice=%s speed=%s",
            workflow_id,
            model_in,
            resolved_model_name,
            provider_type or "<none>",
            body.voice,
            body.speed,
        )
        if provider and provider_type == "qianwen":
            api_key = ((provider or {}).get("api_key") or "").strip()
            if not api_key:
                raise RuntimeError("该模型提供商未配置 API Key")
            base_url = ((provider or {}).get("base_url") or "").strip()
            _qianwen_tts_to_file(
                api_key=api_key,
                base_url=base_url,
                model_name=resolved_model_name,
                text=script,
                voice=body.voice,
                out_file=out,
            )
        else:
            client, model_name = _resolve_tts_client_and_model(model_in)
            speech = client.audio.speech.create(
                model=model_name,
                voice=body.voice,
                input=script[:8000],
                speed=body.speed,
            )
            speech.stream_to_file(str(out))
    except Exception as e:
        logger.exception(
            "[TTS] stage5 failed workflow=%s model=%s voice=%s output=%s",
            workflow_id,
            model_in,
            body.voice,
            body.output_file,
        )
        raise HTTPException(status_code=400, detail=f"TTS 生成失败: {str(e)}")

    _touch_stage(workflow_id, "stage5", "done", outputs=[body.output_file], extra={"bytes": out.stat().st_size})
    return {"ok": True, "workflow_id": workflow_id, "output": body.output_file, "bytes": out.stat().st_size}


@router.post("/workflows/{workflow_id}/stage6/render-video")
async def stage6_render_video(
    workflow_id: str,
    body: Stage6RunIn,
    user: User = Depends(require_teacher),
):
    _ = user
    ppt = _safe_workflow_file(workflow_id, body.ppt_file)
    if not ppt.exists():
        ppt = _safe_workflow_file(workflow_id, body.fallback_ppt)
    if not ppt.exists():
        raise HTTPException(status_code=404, detail="未找到 PPT 文件")

    audio = _safe_workflow_file(workflow_id, body.audio_file)
    if not audio.exists():
        raise HTTPException(status_code=404, detail="未找到音频文件")

    durations: list[float] | None = None
    if body.timing_file:
        timeline = _safe_workflow_file(workflow_id, body.timing_file)
        if timeline.exists():
            data = json.loads(timeline.read_text(encoding="utf-8"))
            if isinstance(data, dict) and isinstance(data.get("durations"), list):
                durations = [max(1.0, float(x)) for x in data["durations"] if isinstance(x, (int, float))]

    # 无时间轴时由 _build_stage6_video 按「每页字数/总字数*音频时长」或均分计算每页时长
    if durations is None:
        pass  # 保持 None

    script_segments_path: Path | None = None
    if body.script_segments_file:
        seg_path = _safe_workflow_file(workflow_id, body.script_segments_file)
        if seg_path.exists():
            script_segments_path = seg_path

    out = _safe_workflow_file(workflow_id, body.output_file)
    info = _build_stage6_video(
        ppt_path=ppt,
        audio_path=audio,
        out_video=out,
        durations=durations,
        script_segments_path=script_segments_path,
    )

    _touch_stage(workflow_id, "stage6", "done", outputs=[body.output_file], extra=info)
    return {"ok": True, "workflow_id": workflow_id, "output": body.output_file, **info}


@router.get("/workflows/{workflow_id}/files/download")
async def download_file(
    workflow_id: str,
    path: str = Query(..., description="工作流内相对路径"),
    user: User = Depends(require_teacher),
):
    _ = user
    file_path = _safe_workflow_file(workflow_id, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="文件不存在")
    return FileResponse(path=file_path, filename=file_path.name)


@router.post("/workflows/{workflow_id}/files/upload")
async def upload_file_override(
    workflow_id: str,
    path: str = Form(..., description="工作流内相对路径"),
    file: UploadFile = File(...),
    user: User = Depends(require_teacher),
):
    _ = user
    dst = _safe_workflow_file(workflow_id, path)
    dst.parent.mkdir(parents=True, exist_ok=True)
    data = await file.read()
    dst.write_bytes(data)
    return {"ok": True, "workflow_id": workflow_id, "path": path, "bytes": len(data)}


@router.get("/workflows/{workflow_id}/files/text", response_model=FileTextOut)
async def read_workflow_text_file(
    workflow_id: str,
    path: str = Query(..., description="工作流内相对路径"),
    user: User = Depends(require_teacher),
):
    _ = user
    file_path = _safe_workflow_file(workflow_id, path)
    if not file_path.exists() or not file_path.is_file():
        raise HTTPException(status_code=404, detail="文件不存在")
    try:
        content = file_path.read_text(encoding="utf-8")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"文本读取失败: {str(e)}")
    return FileTextOut(path=path, content=content)


@router.put("/workflows/{workflow_id}/files/text", response_model=FileTextOut)
async def save_workflow_text_file(
    workflow_id: str,
    body: FileTextIn,
    user: User = Depends(require_teacher),
):
    _ = user
    path = body.path.strip()
    if not path:
        raise HTTPException(status_code=400, detail="path 不能为空")
    dst = _safe_workflow_file(workflow_id, path)
    dst.parent.mkdir(parents=True, exist_ok=True)
    if path.lower().endswith(".json"):
        try:
            json.loads(body.content or "{}")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"JSON 格式不合法: {str(e)}")
    dst.write_text(body.content, encoding="utf-8")
    return FileTextOut(path=path, content=body.content)
